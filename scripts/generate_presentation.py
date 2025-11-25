#!/usr/bin/env python3
"""
Generate a PowerPoint presentation summarizing the Crypto Volatility Detection project.
For 4 presenters based on team charter roles.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BG = RGBColor(0x1a, 0x1a, 0x2e)
ACCENT_BLUE = RGBColor(0x00, 0xd4, 0xff)
ACCENT_GREEN = RGBColor(0x00, 0xff, 0x88)
ACCENT_PURPLE = RGBColor(0x8b, 0x5c, 0xf6)
WHITE = RGBColor(0xff, 0xff, 0xff)
GRAY = RGBColor(0xaa, 0xaa, 0xaa)


def add_title_slide(prs, title, subtitle, presenter):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Presenter
    pres_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.5))
    tf = pres_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Presented by: {presenter}"
    p.font.size = Pt(20)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_content_slide(prs, title, content_items, presenter_role=None):
    """Add a content slide with bullet points."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(0x2a, 0x2a, 0x4a)
    title_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        if isinstance(item, tuple):
            # (text, level) format
            text, level = item
            p.text = "• " + text if level == 0 else "    ‣ " + text
            p.font.size = Pt(22) if level == 0 else Pt(18)
        else:
            p.text = "• " + item
            p.font.size = Pt(22)
        
        p.font.color.rgb = WHITE
        p.space_after = Pt(12)
    
    # Presenter role badge
    if presenter_role:
        badge = slide.shapes.add_textbox(Inches(10.5), Inches(6.8), Inches(2.5), Inches(0.4))
        tf = badge.text_frame
        p = tf.paragraphs[0]
        p.text = presenter_role
        p.font.size = Pt(14)
        p.font.color.rgb = ACCENT_PURPLE
        p.alignment = PP_ALIGN.RIGHT
    
    return slide


def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items, presenter_role=None):
    """Add a two-column content slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(0x2a, 0x2a, 0x4a)
    title_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(5.8), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    # Left column content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(5.8), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(8)
    
    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(5.8), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    
    # Right column content
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(2), Inches(5.8), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(8)
    
    # Presenter role badge
    if presenter_role:
        badge = slide.shapes.add_textbox(Inches(10.5), Inches(6.8), Inches(2.5), Inches(0.4))
        tf = badge.text_frame
        p = tf.paragraphs[0]
        p.text = presenter_role
        p.font.size = Pt(14)
        p.font.color.rgb = ACCENT_PURPLE
        p.alignment = PP_ALIGN.RIGHT
    
    return slide


def add_metrics_slide(prs, title, metrics, presenter_role=None):
    """Add a slide with key metrics in boxes."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(0x2a, 0x2a, 0x4a)
    title_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Metric boxes
    box_width = Inches(2.8)
    box_height = Inches(2)
    start_x = Inches(0.7)
    start_y = Inches(1.8)
    gap = Inches(0.3)
    
    colors = [ACCENT_BLUE, ACCENT_GREEN, ACCENT_PURPLE, RGBColor(0xff, 0x88, 0x00)]
    
    for i, (label, value) in enumerate(metrics):
        row = i // 4
        col = i % 4
        x = start_x + col * (box_width + gap)
        y = start_y + row * (box_height + gap)
        
        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_width, box_height)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x2a, 0x2a, 0x4a)
        box.line.color.rgb = colors[i % len(colors)]
        box.line.width = Pt(2)
        
        # Value
        val_box = slide.shapes.add_textbox(x, y + Inches(0.3), box_width, Inches(1))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(value)
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = colors[i % len(colors)]
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        lbl_box = slide.shapes.add_textbox(x, y + Inches(1.3), box_width, Inches(0.5))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER
    
    # Presenter role badge
    if presenter_role:
        badge = slide.shapes.add_textbox(Inches(10.5), Inches(6.8), Inches(2.5), Inches(0.4))
        tf = badge.text_frame
        p = tf.paragraphs[0]
        p.text = presenter_role
        p.font.size = Pt(14)
        p.font.color.rgb = ACCENT_PURPLE
        p.alignment = PP_ALIGN.RIGHT
    
    return slide


# ============================================================================
# SLIDE 1: Title Slide (Team Lead)
# ============================================================================
add_title_slide(
    prs,
    "Crypto Volatility Detection",
    "Real-Time AI Service for BTC-USD Spike Prediction",
    "Team Lead"
)

# ============================================================================
# SLIDE 2: Project Overview (Team Lead)
# ============================================================================
add_content_slide(
    prs,
    "Project Overview",
    [
        "Objective: Predict short-term volatility spikes in BTC-USD markets",
        ("Binary classification: Spike (1) vs Normal (0)", 1),
        ("60-second lookahead prediction window", 1),
        "Real-time streaming from Coinbase WebSocket API",
        ("Process 1+ tick/second with sub-800ms latency", 1),
        "Full MLOps pipeline with monitoring and rollback",
        ("Docker-based deployment with one-command startup", 1),
        ("Grafana dashboards for real-time observability", 1),
    ],
    "Team Lead"
)

# ============================================================================
# SLIDE 3: Architecture (Team Lead)
# ============================================================================
add_two_column_slide(
    prs,
    "System Architecture",
    "Data Flow",
    [
        "Coinbase WebSocket → Ingestor",
        "Ingestor → Kafka (KRaft mode)",
        "Kafka → Feature Engineering",
        "Features → FastAPI /predict",
        "API → Prometheus metrics",
        "Prometheus → Grafana dashboards",
    ],
    "Tech Stack",
    [
        "API: FastAPI (Python 3.11)",
        "Streaming: Apache Kafka",
        "ML: scikit-learn (Logistic Regression)",
        "Tracking: MLflow",
        "Monitoring: Prometheus + Grafana",
        "Container: Docker Compose",
    ],
    "Team Lead"
)

# ============================================================================
# SLIDE 4: Model Selection (ML Engineer)
# ============================================================================
add_two_column_slide(
    prs,
    "Model Selection & Performance",
    "Models Evaluated",
    [
        "Z-Score Baseline: PR-AUC 0.33",
        "Random Forest: PR-AUC 0.30",
        "Gradient Boosting: PR-AUC 0.84",
        "Logistic Regression: PR-AUC 0.89 ✓",
        "Selected: Logistic Regression",
        "GridSearchCV hyperparameter tuning",
    ],
    "Best Model Metrics",
    [
        "PR-AUC: 0.8917",
        "ROC-AUC: 0.9399",
        "F1 Score: 0.9091",
        "Precision: 94.34%",
        "Recall: 87.72%",
        "Inference: 0.003 ms/sample",
    ],
    "ML Engineer"
)

# ============================================================================
# SLIDE 5: Features & Training (ML Engineer)
# ============================================================================
add_two_column_slide(
    prs,
    "Feature Engineering & Training",
    "Top Features (by importance)",
    [
        "realized_volatility_300s (34%)",
        "realized_volatility_60s (29%)",
        "log_return_300s (16%)",
        "spread_mean_300s (9%)",
        "trade_intensity_300s (7%)",
        "order_book_imbalance (5%)",
    ],
    "Training Details",
    [
        "Dataset: 1,140 samples",
        "Train/Test Split: 80/20",
        "Class Balance: 57% normal, 43% spike",
        "5-Fold Cross Validation",
        "GridSearchCV optimization",
        "All CPU cores utilized (n_jobs=-1)",
    ],
    "ML Engineer"
)

# ============================================================================
# SLIDE 6: API & Infrastructure (Backend/DevOps)
# ============================================================================
add_two_column_slide(
    prs,
    "API & Infrastructure",
    "FastAPI Endpoints",
    [
        "POST /predict - Volatility prediction",
        "GET /health - Service health check",
        "GET /version - Model version info",
        "GET /metrics - Prometheus metrics",
        "Rate limiting: 100 req/min",
        "Structured JSON logging",
    ],
    "Infrastructure",
    [
        "Docker Compose orchestration",
        "Kafka KRaft (no Zookeeper)",
        "Prometheus metrics collection",
        "Grafana dashboards",
        "CI/CD: GitHub Actions",
        "Black + Ruff linting",
    ],
    "Backend/DevOps"
)

# ============================================================================
# SLIDE 7: Monitoring & SLOs (Backend/DevOps)
# ============================================================================
add_metrics_slide(
    prs,
    "Monitoring & SLOs",
    [
        ("p95 Latency Target", "≤800ms"),
        ("Availability", "99.5%"),
        ("Error Rate", "<1%"),
        ("Success Rate", "≥99%"),
        ("CPU Usage", "Real-time"),
        ("Memory", "Real-time"),
        ("Drift Detection", "Evidently"),
        ("Rollback", "ml|baseline"),
    ],
    "Backend/DevOps"
)

# ============================================================================
# SLIDE 8: Demo & Deliverables (All)
# ============================================================================
add_content_slide(
    prs,
    "Demo & Deliverables",
    [
        "One-command startup: docker compose up -d",
        "API Contract: POST /predict with {rows: [...]}",
        ("Returns: {scores, model_variant, version, ts}", 1),
        "Documentation suite:",
        ("team_charter.md - Roles & responsibilities", 1),
        ("selection_rationale.md - Model choice", 1),
        ("slo.md - Service Level Objectives", 1),
        ("runbook.md - Operations guide", 1),
        "Demo video: 8-minute walkthrough",
        ("Startup → Prediction → Failure Recovery → Rollback", 1),
    ],
    "All Team Members"
)

# ============================================================================
# Save the presentation
# ============================================================================
output_path = os.path.join(os.path.dirname(__file__), "..", "docs", "Crypto_Volatility_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")

